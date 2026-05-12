"""
Builds the final-presentation PowerPoint deck for the Austin-San Marcos
Housing Affordability project. Re-run after data/narrative changes:
    python build_pptx.py
Produces:  Austin_SanMarcos_Affordability.pptx
"""
from pathlib import Path

import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).parent
DATA = ROOT / "data"
IMG = ROOT / "_pptx_images"
IMG.mkdir(exist_ok=True)
OUT = ROOT / "Austin_SanMarcos_Affordability.pptx"

ACCENT = RGBColor(0x06, 0xB6, 0xD4)
ACCENT2 = RGBColor(0xEC, 0x48, 0x99)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x47, 0x55, 0x69)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)


# ---------- data loading (mirrors app.py) ----------
def load_zhvi():
    df = pd.read_csv(DATA / "zillow_zhvi_city.csv")
    df = df[df["RegionName"].isin(["Austin", "San Marcos"])].copy()
    date_cols = [c for c in df.columns if c.startswith("20")]
    id_cols = [c for c in df.columns if c not in date_cols]
    out = df.melt(id_vars=id_cols, value_vars=date_cols, var_name="Date", value_name="ZHVI")
    out["Date"] = pd.to_datetime(out["Date"])
    return out


def load_income():
    df = pd.read_csv(DATA / "zillow_income_needed_metro.csv")
    df = df[df["RegionName"].str.contains("Austin", na=False, case=False)].copy()
    date_cols = [c for c in df.columns if c.startswith("20")]
    id_cols = [c for c in df.columns if c not in date_cols]
    out = df.melt(id_vars=id_cols, value_vars=date_cols, var_name="Date", value_name="Income_Needed")
    out["Date"] = pd.to_datetime(out["Date"])
    return out


def load_housing():
    df = pd.read_csv(DATA / "austin_affordable_housing.csv")
    return df.dropna(subset=["Latitude", "Longitude"]).copy()


# ---------- chart rendering ----------
def style(fig, title):
    fig.update_layout(
        title=dict(text=title, font=dict(size=22, color="#0f172a", family="Helvetica")),
        plot_bgcolor="white",
        paper_bgcolor="white",
        font=dict(color="#0f172a", family="Helvetica", size=14),
        margin=dict(l=60, r=30, t=70, b=60),
        legend=dict(bgcolor="rgba(0,0,0,0)"),
    )
    fig.update_xaxes(gridcolor="#e2e8f0", showline=True, linecolor="#cbd5e1")
    fig.update_yaxes(gridcolor="#e2e8f0", showline=True, linecolor="#cbd5e1")


def chart_income(df):
    fig = px.line(df, x="Date", y="Income_Needed", color_discrete_sequence=["#06b6d4"])
    fig.update_traces(line=dict(width=4))
    fig.add_hline(
        y=78000, line_dash="dot", line_color="#94a3b8",
        annotation_text="≈ $78K (2020)", annotation_position="bottom right",
    )
    fig.add_hline(
        y=150000, line_dash="dot", line_color="#ec4899",
        annotation_text="≈ $150K peak (Oct 2022)", annotation_position="top right",
    )
    fig.update_yaxes(title="Income required for median home ($)", tickprefix="$", tickformat=",")
    fig.update_xaxes(title=None)
    style(fig, "The Affordability Wall — Austin Metro")
    return fig


def chart_zhvi(df):
    fig = px.line(
        df, x="Date", y="ZHVI", color="RegionName",
        color_discrete_map={"Austin": "#06b6d4", "San Marcos": "#ec4899"},
    )
    fig.update_traces(line=dict(width=4))
    fig.update_yaxes(title="Typical home value (ZHVI, $)", tickprefix="$", tickformat=",")
    fig.update_xaxes(title=None)
    fig.update_layout(legend_title_text="")
    style(fig, "Synchronized Appreciation Along I-35")
    return fig


def chart_spatial(df):
    df = df.copy()
    units_col = "Affordable_Units" if "Affordable_Units" in df.columns else None
    fig = px.scatter(
        df, x="Longitude", y="Latitude",
        size=units_col, color=units_col,
        color_continuous_scale="Reds",
        size_max=22,
    )
    # I-35 runs roughly along longitude -97.74 in Austin
    fig.add_vline(
        x=-97.74, line_dash="dash", line_color="#0f172a",
        annotation_text="I-35", annotation_position="top",
    )
    fig.update_xaxes(title="Longitude")
    fig.update_yaxes(title="Latitude", scaleanchor="x", scaleratio=1)
    fig.update_layout(coloraxis_colorbar=dict(title="Units"))
    style(fig, "Affordable Housing Concentrates East of I-35")
    return fig


def render_all():
    zhvi = load_zhvi()
    income = load_income()
    housing = load_housing()

    images = {
        "income": IMG / "income_wall.png",
        "zhvi": IMG / "zhvi.png",
        "spatial": IMG / "spatial.png",
    }
    chart_income(income).write_image(images["income"], width=1400, height=750, scale=2)
    chart_zhvi(zhvi).write_image(images["zhvi"], width=1400, height=750, scale=2)
    chart_spatial(housing).write_image(images["spatial"], width=1400, height=900, scale=2)
    return images


# ---------- slide builders ----------
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False,
                color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, items, *, size=20, color=INK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        bullet = p.add_run()
        bullet.text = "•  "
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = ACCENT
        bullet.font.bold = True
        bullet.font.name = "Calibri"
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def add_accent_bar(slide, top=Inches(1.05), width=Inches(0.6)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    return bar


def slide_title(prs):
    s = add_blank(prs)
    # accent block
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SLIDE_H)
    block.fill.solid(); block.fill.fore_color.rgb = ACCENT; block.line.fill.background()
    add_textbox(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.6),
                "Austin–San Marcos\nHousing Affordability",
                size=54, bold=True, color=INK)
    # multi-line trick: replace newline by separate paragraphs
    tb = s.shapes[-1]
    tf = tb.text_frame
    tf.clear()
    for i, line in enumerate(["Austin–San Marcos", "Housing Affordability"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = line
        r.font.size = Pt(54); r.font.bold = True; r.font.color.rgb = INK; r.font.name = "Calibri"

    add_textbox(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.5),
                "A data-driven look at pricing pressure, spatial inequality, and the legacy of zoning.",
                size=20, color=MUTED)
    add_textbox(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.4),
                "Erik Olvera   •   CS 4379G   •   Dr. Aniruddha Bora",
                size=16, color=INK, bold=True)
    add_textbox(s, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.4),
                "Final Presentation — May 2026",
                size=14, color=MUTED)


def slide_section_header(prs, eyebrow, title):
    s = add_blank(prs)
    add_accent_bar(s, top=Inches(0.85))
    add_textbox(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.5),
                eyebrow.upper(), size=14, bold=True, color=ACCENT)
    add_textbox(s, Inches(0.6), Inches(1.2), Inches(12), Inches(1.0),
                title, size=40, bold=True, color=INK)
    return s


def slide_motivation(prs):
    s = slide_section_header(prs, "01 — Motivation", "Why this matters")
    add_bullets(s, Inches(0.7), Inches(2.5), Inches(12), Inches(4.5), [
        "Austin's population grew ~30% in the 2010s — driven by a tech-sector boom.",
        "Median home prices and required incomes have decoupled from local wages.",
        "Spillover demand pressures San Marcos and the I-35 corridor.",
        "Question: how severe is the affordability gap, and where is it concentrated?",
    ], size=22)


def slide_methodology(prs):
    s = slide_section_header(prs, "02 — Data & Methodology", "Four datasets, one story")
    # left column: datasets
    add_textbox(s, Inches(0.7), Inches(2.4), Inches(6), Inches(0.4),
                "DATASETS", size=14, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.7), Inches(2.9), Inches(6.0), Inches(4.0), [
        "Zillow ZHVI — typical home values, Austin & San Marcos",
        "Zillow Income-Needed — required income for median home",
        "Redfin Data Center — weekly inventory & median sale price",
        "City of Austin Open Data — affordable housing locations",
    ], size=18)
    # right column: methods
    add_textbox(s, Inches(7.1), Inches(2.4), Inches(6), Inches(0.4),
                "METHODS", size=14, bold=True, color=ACCENT)
    add_bullets(s, Inches(7.1), Inches(2.9), Inches(5.8), Inches(4.0), [
        "Cleaning: dollar/comma parsing, UTF-16 decoding, NaN drops",
        "Wide → long reshape (pandas melt) for time-series",
        "Plotly visualizations + interactive Streamlit dashboard",
        "Geospatial scatter to surface clustering along I-35",
    ], size=18)


def slide_finding(prs, num, title, image_path, takeaway):
    s = add_blank(prs)
    add_textbox(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.4),
                f"FINDING 0{num}", size=14, bold=True, color=ACCENT)
    add_textbox(s, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
                title, size=30, bold=True, color=INK)
    s.shapes.add_picture(str(image_path), Inches(0.6), Inches(1.85),
                         width=Inches(12.1))
    # Takeaway bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.7),
                             Inches(12.1), Inches(0.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = INK; bar.line.fill.background()
    tb = bar.text_frame
    tb.margin_left = Inches(0.2); tb.margin_right = Inches(0.2)
    p = tb.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Takeaway:  "
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = takeaway
    r2.font.size = Pt(16); r2.font.color.rgb = PAPER; r2.font.name = "Calibri"


def slide_finding3_context(prs):
    """Adds the historical-context slide that pairs with Finding 3."""
    s = slide_section_header(prs, "Historical Context", "The 1928 Austin City Plan")
    add_bullets(s, Inches(0.7), Inches(2.5), Inches(12), Inches(4.5), [
        "1928 plan used redlining and exclusionary zoning to relegate minority and "
        "low-income communities east of I-35.",
        "Decades of underinvestment made East Austin the default site for "
        "city-sponsored affordable housing.",
        "Today's bubble map is the visible footprint of that policy — not a "
        "neutral distribution.",
        "Compounding crisis: as East Austin gentrifies, the populations dependent "
        "on this concentrated infrastructure are being displaced.",
    ], size=20)


def slide_demo(prs):
    s = add_blank(prs)
    # full-bleed dark background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = INK; bg.line.fill.background()
    add_textbox(s, 0, Inches(2.6), SLIDE_W, Inches(0.5),
                "LIVE DEMO", size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, 0, Inches(3.1), SLIDE_W, Inches(1.4),
                "Streamlit Dashboard", size=64, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
    add_textbox(s, 0, Inches(4.7), SLIDE_W, Inches(0.6),
                "ZHVI trends · Affordability wall · Inventory dynamics · Geospatial bubble map",
                size=18, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    add_textbox(s, 0, Inches(6.4), SLIDE_W, Inches(0.4),
                "streamlit run app.py", size=14, color=ACCENT2,
                align=PP_ALIGN.CENTER, font="Menlo")


def slide_limitations(prs):
    s = slide_section_header(prs, "04 — Limitations & Future Work", "What this analysis can't tell us")
    add_textbox(s, Inches(0.7), Inches(2.4), Inches(6), Inches(0.4),
                "LIMITATIONS", size=14, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.7), Inches(2.9), Inches(6.0), Inches(4.0), [
        "Federal interest-rate effects not modeled",
        "Inventory data is metro-aggregated, not by ZIP",
        "Naturally Occurring Affordable Housing (NOAH) is excluded from the city dataset",
    ], size=18)
    add_textbox(s, Inches(7.1), Inches(2.4), Inches(6), Inches(0.4),
                "FUTURE WORK", size=14, bold=True, color=ACCENT)
    add_bullets(s, Inches(7.1), Inches(2.9), Inches(5.8), Inches(4.0), [
        "Forecast when required income reconverges with median wages",
        "Layer in census-tract income data to compute affordability gaps",
        "Add NOAH inventory via parcel-level data scraping",
    ], size=18)


def slide_thanks(prs):
    s = add_blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = PAPER; bg.line.fill.background()
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SLIDE_H)
    block.fill.solid(); block.fill.fore_color.rgb = ACCENT; block.line.fill.background()
    add_textbox(s, Inches(0.9), Inches(2.8), Inches(12), Inches(1.4),
                "Thank you.", size=72, bold=True, color=INK)
    add_textbox(s, Inches(0.9), Inches(4.3), Inches(12), Inches(0.5),
                "Questions?", size=28, color=ACCENT, bold=True)
    add_textbox(s, Inches(0.9), Inches(6.5), Inches(12), Inches(0.4),
                "Erik Olvera   •   CS 4379G   •   Spring 2026",
                size=14, color=MUTED)


# ---------- main ----------
def main():
    images = render_all()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_motivation(prs)
    slide_methodology(prs)
    slide_finding(
        prs, 1,
        "The Affordability Wall",
        images["income"],
        "Required income to afford a median home nearly doubled — from ~$78K (2020) to a ~$150K peak (Oct 2022) — pricing out median earners.",
    )
    slide_finding(
        prs, 2,
        "Synchronized I-35 Appreciation",
        images["zhvi"],
        "Austin and San Marcos move in lockstep — spillover demand inflates both markets proportionally.",
    )
    slide_finding(
        prs, 3,
        "Spatial Inequality of Affordable Housing",
        images["spatial"],
        "Affordable housing clusters east of I-35; West and Central Austin are housing deserts.",
    )
    slide_finding3_context(prs)
    slide_demo(prs)
    slide_limitations(prs)
    slide_thanks(prs)

    prs.save(OUT)
    print(f"Wrote {OUT.relative_to(ROOT)}  ({OUT.stat().st_size // 1024} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
